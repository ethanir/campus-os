from pathlib import Path
import fitz  # PyMuPDF
from pptx import Presentation


def extract_text(file_path: str) -> str:
    """Extract text from a file based on its extension."""
    path = Path(file_path)
    ext = path.suffix.lower()

    extractors = {
        ".pdf": _extract_pdf,
        ".pptx": _extract_pptx,
        ".txt": _extract_plain,
        ".md": _extract_plain,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        return _extract_plain(file_path)

    return extractor(file_path)


def _extract_pdf(file_path: str) -> str:
    """Extract text from PDF using PyMuPDF with superscript detection."""
    doc = fitz.open(file_path)
    pages = []
    for i, page in enumerate(doc):
        text = _extract_page_with_superscripts(page)
        has_images = bool(page.get_images())
        if text.strip():
            header = f"--- PAGE {i + 1}"
            if has_images:
                header += " [CONTAINS FIGURES/IMAGES]"
            header += " ---"
            pages.append(f"{header}\n{text.strip()}")
    doc.close()
    return "\n\n".join(pages)


def _extract_page_with_superscripts(page) -> str:
    """Extract text preserving superscripts/subscripts as ^() and _() notation."""
    blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
    result = []
    for block in blocks:
        if block["type"] != 0:  # skip image blocks
            continue
        for line in block["lines"]:
            spans = line["spans"]
            if not spans:
                continue
            # Find the dominant (most common) font size in this line
            sizes = [s["size"] for s in spans]
            if not sizes:
                continue
            dominant_size = max(set(sizes), key=sizes.count)
            line_text = ""
            for span in spans:
                text = span["text"]
                if not text:
                    continue
                size_ratio = span["size"] / dominant_size if dominant_size > 0 else 1
                # Check vertical position relative to line baseline
                span_y = span["origin"][1]
                line_y = line["spans"][0]["origin"][1]
                y_offset = line_y - span_y  # positive = above baseline
                
                if y_offset > 2:
                    # Superscript: smaller font, positioned above baseline
                    line_text += "^(" + text.strip() + ")"
                elif y_offset < -2:
                    # Subscript: smaller font, positioned below baseline
                    line_text += "_(" + text.strip() + ")"
                else:
                    line_text += text
            result.append(line_text)
    return "\n".join(result)


def _extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint slides."""
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"--- SLIDE {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_plain(file_path: str) -> str:
    """Read plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()
